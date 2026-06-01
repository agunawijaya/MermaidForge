"""
Mermaid to PPTX Generator - v3
================================

Adds to v2:
  * Auto-sized nodes: width/height computed from text length so nothing overflows
  * Clean flat theme (configurable palette): no default PowerPoint gradient blue
  * Theming also applied to cluster borders + labels

Syntax supported (same as v2):
  flowchart TD | LR | BT | RL
  Shapes:       A[Rect]  B(Rounded)  C{Diamond}  D([Pill])  E((Circle))
                F[[Subroutine]]  G[(Database)]  H{{Hexagon}}
                I[/Parallelogram/]  J[\\Reverse\\]
  Edges:        -->  ---  -.->  -.-  ==>  ===  (with optional |label|)
  Grouping:     subgraph Name ... end  (supports direction + nesting)
"""

import re
import subprocess
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ========================================================================
# THEMES
# ========================================================================

THEMES = {
    'clean': {
        'node_fill':       'E8EEF7',   # soft blue
        'node_border':     '2E5C8A',   # medium blue
        'node_border_w':   12700,      # 1pt
        'node_text':       '1F3864',   # dark blue
        'edge_color':      '595959',   # gray
        'edge_label_bg':   'FFFFFF',
        'edge_label_text': '333333',
        'cluster_border':  '9E9E9E',
        'cluster_fill':    None,       # transparent
        'cluster_label':   '555555',
    },
    'dark': {
        'node_fill':       '2B3A55',
        'node_border':     '7EA6E0',
        'node_border_w':   12700,
        'node_text':       'E8EEF7',
        'edge_color':      'CCCCCC',
        'edge_label_bg':   '2B3A55',
        'edge_label_text': 'E8EEF7',
        'cluster_border':  '6C7A96',
        'cluster_fill':    None,
        'cluster_label':   'AAB7C6',
    },
    'mermaid_like': {
        'node_fill':       'ECECFF',
        'node_border':     '9370DB',
        'node_border_w':   12700,
        'node_text':       '333333',
        'edge_color':      '333333',
        'edge_label_bg':   'FFFFFF',
        'edge_label_text': '333333',
        'cluster_border':  'B0B0B0',
        'cluster_fill':    None,
        'cluster_label':   '555555',
    },
}


# ========================================================================
# TEXT-SIZE ESTIMATION (for auto-sized nodes)
# ========================================================================

FONT_PT = 12
FONT_NAME = 'Calibri'
AVG_CHAR_PT = 10.5       # very generous — real render wider than Calibri metrics
LINE_HEIGHT_PT = 18
PAD_H_IN = 0.55
PAD_V_IN = 0.30
MIN_W_IN = 1.8
MIN_H_IN = 0.70
MAX_W_IN = 3.6


def estimate_node_size(text, shape_kind):
    """Return (width_in, height_in) sized to fit text without truncation."""
    text = text or ''
    pt_per_in = 72.0
    text_w_in = len(text) * AVG_CHAR_PT / pt_per_in + PAD_H_IN
    text_h_in = LINE_HEIGHT_PT / pt_per_in + PAD_V_IN

    if text_w_in > MAX_W_IN:
        target_w = MAX_W_IN
        chars_per_line = max(1, int((target_w - PAD_H_IN) * pt_per_in / AVG_CHAR_PT))
        lines = max(1, -(-len(text) // chars_per_line))
        text_w_in = target_w
        text_h_in = (lines * LINE_HEIGHT_PT) / pt_per_in + PAD_V_IN

    w = max(MIN_W_IN, text_w_in)
    h = max(MIN_H_IN, text_h_in)

    # Geometry-aware extra space — PowerPoint restricts text to the shape's
    # INSCRIBED rectangle (not its bounding box), so we need to oversize
    # the bounding box to give the inscribed rect enough room.
    if shape_kind == 'rhombus':
        # Diamond inscribed rectangle ~ 0.5 W x 0.5 H → double the bbox
        w = max(w * 1.9, w + 1.2)
        h = max(h * 1.9, h + 0.6)
    elif shape_kind == 'hexagon':
        w = w + 0.9      # slanted sides eat ~25% width each end
    elif shape_kind in ('parallelogram', 'parallelogram_r'):
        w = w + 0.6
    elif shape_kind == 'stadium':
        w = w + 0.45
    elif shape_kind == 'subroutine':
        w = w + 0.55     # double vertical bars eat both sides
    elif shape_kind == 'circle':
        # Circle inscribed rect ~ 0.71 diameter → ~1.4x oversize
        d = max(w * 1.5, h * 2.2)
        w = h = d
    elif shape_kind == 'cylinder':
        h = h + 0.2

    return round(w, 3), round(h, 3)


# ========================================================================
# 1. PARSER — structural (hierarchical subgraphs)
# ========================================================================

# Shape patterns (checked in order — longer/more-specific first)
SHAPE_TOKENS = [
    (r'\[\(([^)]+)\)\]',   'cylinder'),    # [(text)]
    (r'\[\[([^\]]+)\]\]',  'subroutine'),  # [[text]]
    (r'\(\[([^\]]+)\]\)',  'stadium'),     # ([text])
    (r'\(\(([^)]+)\)\)',   'circle'),      # ((text))
    (r'\{\{([^}]+)\}\}',   'hexagon'),     # {{text}}
    (r'\[/([^/]+)/\]',     'parallelogram'),     # [/text/]
    (r'\[\\([^\\]+)\\\]',  'parallelogram_r'),   # [\text\]
    (r'\[([^\]]+)\]',      'rect'),        # [text]
    (r'\(([^)]+)\)',       'round'),       # (text)
    (r'\{([^}]+)\}',       'rhombus'),     # {text}
]
NODE_FULL_RE = re.compile(
    r'([A-Za-z_][A-Za-z0-9_]*)\s*(?:' +
    '|'.join(f'({p})' for p, _ in SHAPE_TOKENS) +
    r')'
)

EDGE_RE = re.compile(
    r'([A-Za-z_][A-Za-z0-9_]*)\s*'
    r'(-->|---|-\.->|-\.-|==>|===)'
    r'\s*(?:\|([^|]*)\|\s*)?'
    r'([A-Za-z_][A-Za-z0-9_]*)'
)
EDGE_STYLE = {
    '-->':  ('solid',  True),
    '---':  ('solid',  False),
    '-.->': ('dashed', True),
    '-.-':  ('dashed', False),
    '==>':  ('thick',  True),
    '===':  ('thick',  False),
}


def _strip_shapes_from_line(line):
    """Remove all shape bracket contents (but keep node ids) for edge parsing."""
    out = line
    for pat, _ in SHAPE_TOKENS:
        out = re.sub(pat, '', out)
    return out


def parse_mermaid(code):
    """Parse mermaid flowchart into a structural model.

    Returns:
      nodes: dict {id -> {'text', 'shape', 'parent'}}
      edges: list of {'src', 'dst', 'label', 'style', 'arrow'}
      subgraphs: dict {sid -> {'label', 'parent', 'direction', 'children'}}
                 (children is list of node ids + nested subgraph ids, in order)
      direction: root direction (TB/LR/BT/RL)
    """
    nodes = {}
    edges = []
    subgraphs = {}           # sid -> info
    direction = 'TB'

    subgraph_counter = [0]   # for auto-naming anonymous subgraphs

    # Stack of subgraph ids we're currently inside (root = None)
    sg_stack = [None]

    raw_lines = [l.rstrip() for l in code.strip().split('\n')]

    for line in raw_lines:
        s = line.strip()
        if not s or s.startswith('%%'):
            continue

        m = re.match(r'(?:flowchart|graph)\s+(TB|TD|LR|BT|RL)\b', s, re.I)
        if m:
            d = m.group(1).upper()
            direction = 'TB' if d == 'TD' else d
            continue

        # subgraph open: `subgraph ID` or `subgraph ID [Display Label]` or `subgraph "Label"`
        m = re.match(r'subgraph\s+(.+)$', s)
        if m:
            rest = m.group(1).strip()
            # Case 1: `ID [Label]`
            m2 = re.match(r'([A-Za-z_][A-Za-z0-9_]*)\s*\[([^\]]+)\]\s*$', rest)
            # Case 2: `ID` alone
            m3 = re.match(r'([A-Za-z_][A-Za-z0-9_]*)\s*$', rest)
            # Case 3: quoted label, no id
            m4 = re.match(r'"([^"]+)"\s*$', rest)
            if m2:
                sid, label = m2.group(1), m2.group(2).strip()
            elif m3:
                sid = m3.group(1); label = sid
            elif m4:
                subgraph_counter[0] += 1
                sid = f'_sg{subgraph_counter[0]}'
                label = m4.group(1).strip()
            else:
                subgraph_counter[0] += 1
                sid = f'_sg{subgraph_counter[0]}'
                label = rest
            parent = sg_stack[-1]
            subgraphs[sid] = {'label': label, 'parent': parent,
                              'direction': None, 'children': []}
            if parent is not None:
                subgraphs[parent]['children'].append(sid)
            sg_stack.append(sid)
            continue

        if s == 'end':
            if len(sg_stack) > 1:
                sg_stack.pop()
            continue

        # direction override inside subgraph
        m = re.match(r'direction\s+(TB|TD|LR|BT|RL)\b', s, re.I)
        if m and sg_stack[-1] is not None:
            d = m.group(1).upper()
            subgraphs[sg_stack[-1]]['direction'] = 'TB' if d == 'TD' else d
            continue

        # Normal content line: could be node decls, edges, or both on same line
        current_parent = sg_stack[-1]

        # Extract explicit node shape declarations.
        # Mermaid rule: a node belongs to the subgraph where its SHAPE is declared,
        # not where it's first seen as an edge endpoint.
        for m in NODE_FULL_RE.finditer(s):
            nid = m.group(1)
            for i, (_, kind) in enumerate(SHAPE_TOKENS):
                outer = m.group(2 + 2 * i)
                text = m.group(3 + 2 * i)
                if outer is not None:
                    if nid not in nodes:
                        nodes[nid] = {'text': text.strip(), 'shape': kind,
                                      'parent': current_parent}
                        if current_parent is not None:
                            subgraphs[current_parent]['children'].append(nid)
                    else:
                        # Update shape/text
                        nodes[nid]['text'] = text.strip()
                        nodes[nid]['shape'] = kind
                        # If currently unparented but being declared inside a subgraph,
                        # move it into that subgraph.
                        old_parent = nodes[nid]['parent']
                        if old_parent is None and current_parent is not None:
                            nodes[nid]['parent'] = current_parent
                            subgraphs[current_parent]['children'].append(nid)
                    break

        # Extract edges
        stripped = _strip_shapes_from_line(s)
        pos = 0
        while True:
            m = EDGE_RE.search(stripped, pos)
            if not m:
                break
            src, edge_tok, label, dst = m.group(1), m.group(2), m.group(3) or '', m.group(4)
            for nid in (src, dst):
                if nid not in nodes:
                    nodes[nid] = {'text': nid, 'shape': 'rect',
                                  'parent': current_parent}
                    if current_parent is not None:
                        subgraphs[current_parent]['children'].append(nid)
            style, has_arrow = EDGE_STYLE[edge_tok]
            edges.append({'src': src, 'dst': dst, 'label': label,
                          'style': style, 'arrow': has_arrow})
            pos = m.end()

    # Also reference-only subgraphs that are undeclared root (no parent)
    for sid, info in subgraphs.items():
        if info['parent'] is None:
            pass  # already handled

    return nodes, edges, subgraphs, direction


# ========================================================================
# 2. LAYOUT via Graphviz (with clusters)
# ========================================================================

GV_RANKDIR = {'TB': 'TB', 'LR': 'LR', 'BT': 'BT', 'RL': 'RL'}
PT_PER_INCH = 72.0


def _gv_escape(s):
    return s.replace('\\', '\\\\').replace('"', r'\"')


def _emit_dot(nodes, edges, subgraphs, direction,
              node_sizes, rank_sep_in, node_sep_in):
    """node_sizes: dict {id -> (w_in, h_in)}."""
    lines = [
        f'digraph G {{',
        f'  rankdir={GV_RANKDIR[direction]};',
        f'  graph [ranksep={rank_sep_in}, nodesep={node_sep_in}, compound=true];',
        f'  node [shape=box, fixedsize=true];',
    ]

    def emit_node(pad, nid):
        w, h = node_sizes[nid]
        label = _gv_escape(nodes[nid]['text'])
        lines.append(f'{pad}{nid} [label="{label}", width={w}, height={h}];')

    def emit_subgraph(sid, indent):
        info = subgraphs[sid]
        pad = '  ' * indent
        lines.append(f'{pad}subgraph cluster_{sid} {{')
        lines.append(f'{pad}  label="{_gv_escape(info["label"])}";')
        lines.append(f'{pad}  style=rounded;')
        for child in info['children']:
            if child in nodes:
                emit_node(pad + '  ', child)
            elif child in subgraphs:
                emit_subgraph(child, indent + 1)
        lines.append(f'{pad}}}')

    for nid, nd in nodes.items():
        if nd['parent'] is None:
            emit_node('  ', nid)
    for sid, info in subgraphs.items():
        if info['parent'] is None:
            emit_subgraph(sid, 1)
    for e in edges:
        lines.append(f'  {e["src"]} -> {e["dst"]};')
    lines.append('}')
    return '\n'.join(lines)


def compute_layout(nodes, edges, subgraphs, direction,
                   node_sizes, rank_sep_in=0.7, node_sep_in=0.45):
    """Run Graphviz `dot` with per-node sizes. Returns node_positions,
    cluster_positions, layout_w_in, layout_h_in.
    """
    dot_src = _emit_dot(nodes, edges, subgraphs, direction,
                        node_sizes, rank_sep_in, node_sep_in)
    proc = subprocess.run(['dot', '-Tjson'], input=dot_src,
                          capture_output=True, text=True, check=True)
    data = json.loads(proc.stdout)

    bb = list(map(float, data.get('bb', '0,0,0,0').split(',')))
    bb_x_min, bb_y_min, bb_x_max, bb_y_max = bb
    layout_w = (bb_x_max - bb_x_min) / PT_PER_INCH
    layout_h = (bb_y_max - bb_y_min) / PT_PER_INCH

    def to_slide(x_pt, y_pt):
        return ((x_pt - bb_x_min) / PT_PER_INCH,
                (bb_y_max - y_pt) / PT_PER_INCH)

    node_positions = {}
    cluster_positions = {}

    for obj in data.get('objects', []):
        name = obj.get('name', '')
        if name in nodes:
            x_pt, y_pt = map(float, obj['pos'].split(','))
            w_in = float(obj['width'])
            h_in = float(obj['height'])
            cx, cy = to_slide(x_pt, y_pt)
            node_positions[name] = (cx, cy, w_in, h_in)
        elif name.startswith('cluster_'):
            sid = name[len('cluster_'):]
            if sid in subgraphs and 'bb' in obj:
                cx_min, cy_min, cx_max, cy_max = map(float, obj['bb'].split(','))
                left, top = to_slide(cx_min, cy_max)
                right, bottom = to_slide(cx_max, cy_min)
                cluster_positions[sid] = (left, top, right - left, bottom - top)

    return node_positions, cluster_positions, layout_w, layout_h


# ========================================================================
# 3. PPTX RENDERER
# ========================================================================

SHAPE_MAP = {
    'rect':           MSO_SHAPE.RECTANGLE,
    'round':          MSO_SHAPE.ROUNDED_RECTANGLE,
    'rhombus':        MSO_SHAPE.DIAMOND,
    'stadium':        MSO_SHAPE.ROUNDED_RECTANGLE,
    'circle':         MSO_SHAPE.OVAL,
    'cylinder':       MSO_SHAPE.FLOWCHART_MAGNETIC_DISK,
    'subroutine':     MSO_SHAPE.FLOWCHART_PREDEFINED_PROCESS,
    'hexagon':        MSO_SHAPE.FLOWCHART_PREPARATION,
    'parallelogram':  MSO_SHAPE.PARALLELOGRAM,
    'parallelogram_r': MSO_SHAPE.PARALLELOGRAM,
}

# Connection point indices for preset rectangle/rounded/diamond shapes.
# IMPORTANT: These are PowerPoint's ACTUAL indices (verified by diffing
# manually-drawn connectors). LibreOffice's elbow-routing fallback can make
# the wrong indices appear to work visually, so do NOT trust visual tests.
#   0 = top, 1 = LEFT, 2 = bottom, 3 = RIGHT
CXN_TOP, CXN_LEFT, CXN_BOTTOM, CXN_RIGHT = 0, 1, 2, 3


def choose_connection_points(src_pos, dst_pos):
    sx, sy, _, _ = src_pos
    dx, dy, _, _ = dst_pos
    ddx = dx - sx
    ddy = dy - sy
    if abs(ddx) > abs(ddy):
        return (CXN_RIGHT, CXN_LEFT) if ddx > 0 else (CXN_LEFT, CXN_RIGHT)
    return (CXN_BOTTOM, CXN_TOP) if ddy > 0 else (CXN_TOP, CXN_BOTTOM)


def _set_solid_fill(spPr, hex_color):
    """Set shape fill to a flat solid color (hex without #)."""
    # Remove any existing fill tags
    for tag in ('a:solidFill', 'a:gradFill', 'a:blipFill',
                'a:noFill', 'a:pattFill'):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    sf = etree.SubElement(spPr, qn('a:solidFill'))
    srgb = etree.SubElement(sf, qn('a:srgbClr'))
    srgb.set('val', hex_color)


def _set_no_fill(spPr):
    for tag in ('a:solidFill', 'a:gradFill', 'a:blipFill',
                'a:noFill', 'a:pattFill'):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    etree.SubElement(spPr, qn('a:noFill'))


def _set_line(spPr, hex_color, width_emu, dashed=False):
    for el in spPr.findall(qn('a:ln')):
        spPr.remove(el)
    ln = etree.SubElement(spPr, qn('a:ln'))
    ln.set('w', str(width_emu))
    sf = etree.SubElement(ln, qn('a:solidFill'))
    srgb = etree.SubElement(sf, qn('a:srgbClr'))
    srgb.set('val', hex_color)
    if dashed:
        prstDash = etree.SubElement(ln, qn('a:prstDash'))
        prstDash.set('val', 'dash')


def _disable_style_inheritance(shape_element):
    """Remove <p:style> so the shape uses its own fill/line, not theme defaults.
    Without this, python-pptx's default <p:style> overlays theme accent colors.
    """
    style = shape_element.find(qn('p:style'))
    if style is not None:
        shape_element.remove(style)


def _apply_node_theme(shape, theme):
    sp = shape._element
    _disable_style_inheritance(sp)
    spPr = sp.find(qn('p:spPr'))
    _set_solid_fill(spPr, theme['node_fill'])
    _set_line(spPr, theme['node_border'], theme['node_border_w'])


def _set_shape_text(shape, text, theme, font_size_pt=FONT_PT, wrap='square'):
    """wrap='square' (wrap to bbox) or 'none' (no wrap, overflow allowed)."""
    tf = shape.text_frame
    tf.word_wrap = (wrap == 'square')
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = text
    # Also set wrap attribute explicitly on bodyPr
    bodyPr = shape.text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('wrap', wrap)
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(font_size_pt)
            run.font.name = FONT_NAME
            run.font.color.rgb = RGBColor.from_string(theme['node_text'])


# Shapes where PowerPoint's preset text rectangle is narrower than the bbox —
# these wrap text prematurely even when the bbox is plenty wide.
# For these, disable word-wrap and let text overflow horizontally instead.
NO_WRAP_SHAPES = {'subroutine', 'hexagon', 'parallelogram', 'parallelogram_r'}


def _apply_edge_style(connector, style, has_arrow, theme):
    sp = connector._element
    _disable_style_inheritance(sp)
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    width_emu = {'solid': 12700, 'dashed': 12700, 'thick': 28575}[style]
    _set_line(spPr, theme['edge_color'], width_emu, dashed=(style == 'dashed'))
    if has_arrow:
        ln = spPr.find(qn('a:ln'))
        tailEnd = etree.SubElement(ln, qn('a:tailEnd'))
        tailEnd.set('type', 'triangle')
        tailEnd.set('w', 'med')
        tailEnd.set('h', 'med')


def _add_connector_label(slide, connector, label, theme, font_size_pt=10):
    if not label:
        return
    x1, y1 = connector.begin_x, connector.begin_y
    x2, y2 = connector.end_x, connector.end_y
    mid_x, mid_y = (x1 + x2) // 2, (y1 + y2) // 2
    # Size label to text
    w_in = max(0.5, min(2.0, len(label) * 0.09 + 0.2))
    tb_w, tb_h = Inches(w_in), Inches(0.28)
    tb = slide.shapes.add_textbox(mid_x - tb_w // 2, mid_y - tb_h // 2, tb_w, tb_h)
    # Opaque background so line doesn't show through the text
    sp = tb._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    _set_solid_fill(spPr, theme['edge_label_bg'])
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.text = label
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = RGBColor.from_string(theme['edge_label_text'])


def _add_cluster_background(slide, left_in, top_in, w_in, h_in, label, theme):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in), Inches(top_in),
        Inches(w_in), Inches(h_in)
    )
    sp = shape._element
    _disable_style_inheritance(sp)
    spPr = sp.find(qn('p:spPr'))
    if theme.get('cluster_fill'):
        _set_solid_fill(spPr, theme['cluster_fill'])
    else:
        _set_no_fill(spPr)
    _set_line(spPr, theme['cluster_border'], 12700, dashed=True)

    # Label
    lb_w, lb_h = Inches(min(w_in - 0.2, 3)), Inches(0.35)
    lb = slide.shapes.add_textbox(
        Inches(left_in + 0.15), Inches(top_in + 0.05), lb_w, lb_h)
    tf = lb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.text = label
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.LEFT
        for run in para.runs:
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = RGBColor.from_string(theme['cluster_label'])
    return shape


def generate_pptx(mermaid_code, output_path,
                  slide_w_in=13.333, slide_h_in=7.5,
                  margin_in=0.5,
                  theme='clean'):
    """Main entry. Convert mermaid_code to a .pptx at output_path.
    `theme` can be a key of THEMES or a dict matching its schema.
    """
    nodes, edges, subgraphs, direction = parse_mermaid(mermaid_code)
    if not nodes:
        raise ValueError("No nodes found in Mermaid code.")

    theme_cfg = THEMES[theme] if isinstance(theme, str) else theme

    # Per-node size from text length + shape
    node_sizes = {nid: estimate_node_size(nd['text'], nd['shape'])
                  for nid, nd in nodes.items()}

    node_positions, cluster_positions, layout_w, layout_h = compute_layout(
        nodes, edges, subgraphs, direction, node_sizes)

    # Scale/center into slide
    avail_w = slide_w_in - 2 * margin_in
    avail_h = slide_h_in - 2 * margin_in
    scale = min(avail_w / layout_w, avail_h / layout_h, 1.0) if layout_w and layout_h else 1.0
    off_x = margin_in + (avail_w - layout_w * scale) / 2
    off_y = margin_in + (avail_h - layout_h * scale) / 2

    def xform_center(cx, cy, w, h):
        return (off_x + cx * scale, off_y + cy * scale, w * scale, h * scale)

    def xform_bbox(l, t, w, h):
        return (off_x + l * scale, off_y + t * scale, w * scale, h * scale)

    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 1) Cluster backgrounds first (behind everything)
    for sid, (l, t, w, h) in cluster_positions.items():
        sl, st, sw, sh = xform_bbox(l, t, w, h)
        _add_cluster_background(slide, sl, st, sw, sh,
                                subgraphs[sid]['label'], theme_cfg)

    # 2) Nodes
    shape_refs = {}
    for nid, (cx, cy, w, h) in node_positions.items():
        scx, scy, sw, sh = xform_center(cx, cy, w, h)
        left, top = scx - sw / 2, scy - sh / 2
        mso_shape = SHAPE_MAP.get(nodes[nid]['shape'], MSO_SHAPE.RECTANGLE)
        s = slide.shapes.add_shape(
            mso_shape, Inches(left), Inches(top), Inches(sw), Inches(sh))
        if nodes[nid]['shape'] == 'stadium':
            try:
                s.adjustments[0] = 0.5
            except Exception:
                pass
        _apply_node_theme(s, theme_cfg)
        wrap_mode = 'none' if nodes[nid]['shape'] in NO_WRAP_SHAPES else 'square'
        _set_shape_text(s, nodes[nid]['text'], theme_cfg, wrap=wrap_mode)
        shape_refs[nid] = s

    # 3) Connectors
    scaled_centers = {nid: xform_center(*pos) for nid, pos in node_positions.items()}
    for edge in edges:
        src_c = scaled_centers[edge['src']]
        dst_c = scaled_centers[edge['dst']]
        src_idx, dst_idx = choose_connection_points(src_c, dst_c)
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.ELBOW,
            Inches(src_c[0]), Inches(src_c[1]),
            Inches(dst_c[0]), Inches(dst_c[1])
        )
        conn.begin_connect(shape_refs[edge['src']], src_idx)
        conn.end_connect(shape_refs[edge['dst']], dst_idx)
        _apply_edge_style(conn, edge['style'], edge['arrow'], theme_cfg)
        _add_connector_label(slide, conn, edge['label'], theme_cfg)

    prs.save(output_path)
    return output_path


# ========================================================================
# DEMO
# ========================================================================
if __name__ == '__main__':
    demos = {
        'v3_flowchart': """
        flowchart TD
            Start([Start]) --> Fetch[Fetch user data]
            Fetch --> Cache{Cache hit?}
            Cache -->|Yes| Serve[Serve cached]
            Cache -->|No| DB[(Query PostgreSQL)]
            DB --> Store[Store in cache]
            Store --> Serve
            Serve --> Log[Log request]
            Log --> End([End])
        """,
        'v3_subgraph': """
        flowchart TB
            Start([Start]) --> Login
            subgraph Auth [Authentication]
                Login[Login Page] --> Verify{Credentials OK?}
                Verify -->|Yes| Session[Create session]
            end
            Session --> Dashboard[Dashboard]
            Verify -->|No| Start
        """,
        'v3_nested': """
        flowchart LR
            Req[API Request] --> Gateway
            subgraph Backend [Backend Services]
                Gateway[API Gateway] --> Service
                subgraph DataLayer [Data Layer]
                    Service[User Service] --> Cache[(Redis Cache)]
                    Service --> DB[(PostgreSQL)]
                end
            end
            Cache --> Response[Response]
            DB --> Response
        """,
        'v3_all_shapes': """
        flowchart TD
            A[Rect] --> B(Rounded)
            B --> C{Decision}
            C --> D([Stadium])
            D --> E[[Subroutine]]
            E --> F[(Database)]
            F --> G{{Hexagon}}
            G --> H[/Parallelogram/]
            H --> I((Circle))
        """,
        'v3_dark_theme': """
        flowchart LR
            A[Start] --> B{Valid?}
            B -->|Yes| C[Process]
            B -->|No| D[Reject]
            C --> E([Done])
            D --> E
        """,
    }
    import os
    os.makedirs('./out', exist_ok=True)
    for name, code in demos.items():
        path = f'./out/{name}.pptx'
        theme = 'dark' if 'dark' in name else 'clean'
        generate_pptx(code, path, theme=theme)
        print(f'OK {name} (theme={theme}) -> {path}')
